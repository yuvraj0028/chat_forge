import logging
from pathlib import Path

from rest_framework import viewsets, permissions, status
from rest_framework.decorators import action
from rest_framework.response import Response
from rest_framework.views import APIView

from chat_app.throttles import ChatRateThrottle
from django.conf import settings
from projects.models import Project
from .models import Conversation, Message
from .serializers import (
    ConversationSerializer, ConversationListSerializer,
    MessageSerializer, ChatRequestSerializer,
)
from .llm_service import LLMService, LLMError

logger = logging.getLogger('chat')

logger = logging.getLogger('chat')


class ConversationViewSet(viewsets.ModelViewSet):
    serializer_class = ConversationSerializer
    permission_classes = [permissions.IsAuthenticated]

    def get_queryset(self):
        qs = Conversation.objects.select_related('project').filter(user=self.request.user)
        project_id = self.request.query_params.get('project_id')
        if project_id:
            qs = qs.filter(project_id=project_id)
        return qs

    def get_serializer_class(self):
        if self.action == 'list':
            return ConversationListSerializer
        return ConversationSerializer

    def perform_create(self, serializer):
        serializer.save(user=self.request.user)

    @action(detail=True, methods=['get'])
    def messages(self, request, pk=None):
        conversation = self.get_object()
        messages = conversation.messages.all()
        serializer = MessageSerializer(messages, many=True)
        return Response(serializer.data)

    def perform_destroy(self, instance):
        logger.info('Conversation deleted: %s (id=%s)', instance.title, instance.pk)
        instance.delete()


class ChatView(APIView):
    permission_classes = [permissions.IsAuthenticated]
    throttle_classes = [ChatRateThrottle]

    def post(self, request):
        serializer = ChatRequestSerializer(data=request.data)
        if not serializer.is_valid():
            return Response(
                {'error': 'Invalid request.', 'details': serializer.errors},
                status=status.HTTP_400_BAD_REQUEST,
            )
        data = serializer.validated_data

        try:
            project = Project.objects.get(
                id=data['project_id'],
                owner=request.user,
                is_active=True,
            )
        except Project.DoesNotExist:
            return Response(
                {'error': 'Project not found or has been deactivated.'},
                status=status.HTTP_404_NOT_FOUND,
            )

        if data.get('conversation_id'):
            try:
                conversation = Conversation.objects.get(
                    id=data['conversation_id'],
                    user=request.user,
                    project=project,
                )
            except Conversation.DoesNotExist:
                return Response(
                    {'error': 'Conversation not found.'},
                    status=status.HTTP_404_NOT_FOUND,
                )
        else:
            title = data['message'][:80] + ('...' if len(data['message']) > 80 else '')
            conversation = Conversation.objects.create(
                project=project,
                user=request.user,
                title=title,
            )

        Message.objects.create(
            conversation=conversation,
            role='user',
            content=data['message'],
        )

        memory_window = data.get('memory_window') or project.memory_window
        messages_for_llm = self._build_messages(project, conversation, memory_window)

        try:
            llm = LLMService()
            model = getattr(settings, 'GROQ_MODEL', '') or project.model
            assistant_text, token_stats = llm.generate_content(
                model=model,
                messages=messages_for_llm,
                temperature=data['temperature'],
                max_output_tokens=data['max_tokens'],
            )
        except LLMError as e:
            logger.warning('LLM error for user=%s project=%s: %s', request.user.pk, project.pk, e)
            Message.objects.filter(conversation=conversation, role='user').last().delete()
            if not conversation.messages.exists():
                conversation.delete()
            return Response(
                {'error': str(e)},
                status=status.HTTP_400_BAD_REQUEST,
            )
        except Exception:
            logger.exception('Unexpected error in chat for user=%s', request.user.pk)
            return Response(
                {'error': 'Something went wrong. Please try again.'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR,
            )

        if not assistant_text:
            assistant_text = 'I was unable to generate a response. Please try rephrasing your message.'

        Message.objects.create(
            conversation=conversation,
            role='assistant',
            content=assistant_text,
        )

        conversation.save()

        # Generate a smart title for the first exchange
        if not data.get('conversation_id'):
            try:
                title_prompt = (
                    'Generate a short conversation title (max 6 words) for this exchange. '
                    'Reply with ONLY the title, no quotes, no period.\n'
                    f'User: {data["message"][:200]}\n'
                    f'Assistant: {assistant_text[:200]}'
                )
                title_text, _ = llm.generate_content(
                    model=model,
                    messages=[{'role': 'user', 'content': title_prompt}],
                    temperature=0.3,
                    max_output_tokens=20,
                )
                clean_title = title_text.strip().strip('"\'`.')
                if clean_title and len(clean_title) > 3:
                    conversation.title = clean_title[:80]
                    conversation.save(update_fields=['title', 'updated_at'])
            except Exception:
                logger.debug('Title generation failed, using default')

        logger.info(
            'Chat: user=%s project=%s conv=%s model=%s input=~%d output=~%d',
            request.user.pk, project.pk, conversation.pk, model,
            token_stats.get('input_tokens', 0), token_stats.get('output_tokens', 0),
        )

        return Response({
            'conversation_id': conversation.id,
            'message': assistant_text,
            'model': model,
            'memory_window_used': memory_window,
            'token_usage': token_stats,
        })

    def _build_messages(self, project, conversation, memory_window):
        messages = []

        # System prompt
        prompts = project.prompts.filter(is_system_prompt=True, is_active=True)
        if prompts:
            system_content = '\n\n'.join(p.content for p in prompts)
        else:
            system_content = 'You are a helpful assistant.'

        # Attach file contents to system prompt
        files = project.files.all()
        if files.exists():
            file_parts = []
            for f in files:
                content = self._read_file(f)
                if content:
                    file_parts.append(f"[File: {f.original_name}]\n{content}")
            if file_parts:
                system_content += '\n\n## Uploaded Files\n' + '\n\n'.join(file_parts)

        messages.append({'role': 'system', 'content': system_content})

        # Conversation history
        all_messages = list(conversation.messages.order_by('created_at'))
        recent = all_messages[-memory_window:] if len(all_messages) > memory_window else all_messages

        for msg in recent:
            messages.append({'role': msg.role, 'content': msg.content})

        return messages

    def _read_file(self, project_file):
        """Read text content from an uploaded file."""
        if not project_file.file:
            return None

        try:
            file_path = project_file.file.path
        except Exception:
            return None

        filename = project_file.original_name.lower()
        ext = Path(filename).suffix

        try:
            # PDF
            if ext == '.pdf':
                return self._read_pdf(file_path)

            # Word documents
            if ext == '.docx':
                return self._read_docx(file_path)

            # Excel spreadsheets
            if ext == '.xlsx':
                return self._read_xlsx(file_path)

            # PowerPoint
            if ext == '.pptx':
                return self._read_pptx(file_path)

            # Images
            if ext in ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'):
                return self._read_image(file_path, project_file.original_name)

            # Text/code files
            text_extensions = {
                '.txt', '.md', '.csv', '.json', '.xml', '.yaml', '.yml',
                '.py', '.js', '.ts', '.jsx', '.tsx', '.html', '.css',
                '.java', '.c', '.cpp', '.h', '.cs', '.go', '.rs', '.rb',
                '.php', '.sql', '.sh', '.bat', '.ps1', '.r', '.swift',
                '.kt', '.scala', '.toml', '.ini', '.cfg', '.conf',
                '.log', '.env', '.gitignore', '.dockerfile', '.makefile',
            }
            if ext in text_extensions:
                return self._read_text(file_path)

            return None
        except Exception:
            logger.debug('Could not read file: %s', project_file.original_name)
            return None

    def _truncate(self, text, max_chars=100_000):
        if len(text) > max_chars:
            return text[:max_chars] + f'\n[...truncated, {len(text)} chars total]'
        return text

    def _read_text(self, file_path):
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return self._truncate(f.read())

    def _read_pdf(self, file_path):
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        text = '\n'.join(page.extract_text() or '' for page in reader.pages)
        return self._truncate(text) if text.strip() else '[PDF contains no extractable text]'

    def _read_docx(self, file_path):
        import docx
        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(' | '.join(cells))
        text = '\n'.join(paragraphs)
        return self._truncate(text) if text.strip() else '[Document is empty]'

    def _read_xlsx(self, file_path):
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            parts.append(f'[Sheet: {sheet_name}]')
            # First row as header
            header = [str(c) if c is not None else '' for c in rows[0]]
            parts.append(' | '.join(header))
            parts.append('-' * 40)
            for row in rows[1:100]:  # Cap at 100 rows
                cells = [str(c) if c is not None else '' for c in row]
                parts.append(' | '.join(cells))
            if len(rows) > 101:
                parts.append(f'[...{len(rows) - 101} more rows]')
        wb.close()
        text = '\n'.join(parts)
        return self._truncate(text) if text.strip() else '[Spreadsheet is empty]'

    def _read_pptx(self, file_path):
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f'[Slide {i}]']
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_text.append(' | '.join(cells))
            if len(slide_text) > 1:
                parts.append('\n'.join(slide_text))
        text = '\n\n'.join(parts)
        return self._truncate(text) if text.strip() else '[Presentation is empty]'

    def _read_image(self, file_path, filename):
        from PIL import Image
        img = Image.open(file_path)
        width, height = img.size
        fmt = img.format or 'unknown'
        mode = img.mode

        parts = [
            f'[Image: {filename}]',
            f'Format: {fmt}, Size: {width}x{height}px, Color: {mode}',
        ]

        # Extract EXIF metadata if available
        exif = img.getexif() if hasattr(img, 'getexif') else {}
        if exif:
            interesting = {271: 'Camera Make', 272: 'Camera Model', 306: 'Date', 36867: 'Date Taken'}
            for tag_id, label in interesting.items():
                if tag_id in exif:
                    parts.append(f'{label}: {exif[tag_id]}')

        # Try OCR with pytesseract if available
        try:
            import pytesseract
            ocr_text = pytesseract.image_to_string(img)
            if ocr_text.strip():
                parts.append(f'\n[OCR Text]\n{ocr_text.strip()}')
        except Exception:
            parts.append('[OCR not available - install Tesseract for text extraction]')

        return '\n'.join(parts)
